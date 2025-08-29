"""
Generate a PowerPoint from the outline if python-pptx is available.

Usage:
  python docs/make_ppt_if_possible.py
"""
import os, re

OUTLINE_FILE = os.path.join(os.path.dirname(__file__), "User_Guide_Outline.md")
OUTPUT_PPTX = os.path.join(os.path.dirname(__file__), "User_Guide_AI_Assistant.pptx")

def main():
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as e:
        print("python-pptx not installed. Install with `pip install python-pptx` and rerun.")
        return
    prs = Presentation()
    with open(OUTLINE_FILE, "r", encoding="utf-8") as f:
        content = f.read().splitlines()

    # Simple parser: lines that start with '## Slide' start a new slide.
    slide = None
    for line in content:
        if line.startswith("## Slide"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
            title = slide.shapes.title
            title.text = line.replace("## ", "")
            body = slide.shapes.placeholders[1].text_frame
            body.clear()
        elif line.startswith("# "):
            # Ignore the H1 outline title
            continue
        elif slide is not None:
            if line.strip() == "":
                continue
            if line.startswith("- "):
                body.add_paragraph().text = line[2:]
            else:
                if not body.text:
                    body.text = line
                else:
                    body.add_paragraph().text = line
    prs.save(OUTPUT_PPTX)
    print("Saved:", OUTPUT_PPTX)

if __name__ == "__main__":
    main()
